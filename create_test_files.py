"""Generate 1000-page test files for performance testing."""

from pathlib import Path

# Add reportlab for PDF generation
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.lib import colors

# Add python-pptx for PPTX generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_pdf(filename: str, num_pages: int, variant: str = "A"):
    """Create a PDF with specified number of pages.

    Args:
        filename: Output filename
        num_pages: Number of pages to create
        variant: "A" or "B" - affects colors and some content
    """
    print(f"Creating {filename} with {num_pages} pages...")

    c = canvas.Canvas(filename, pagesize=landscape(A4))
    width, height = landscape(A4)

    # Different colors for A and B variants
    if variant == "A":
        bg_color = colors.HexColor("#E3F2FD")  # Light blue
        accent_color = colors.HexColor("#1976D2")  # Blue
    else:
        bg_color = colors.HexColor("#FFF3E0")  # Light orange
        accent_color = colors.HexColor("#F57C00")  # Orange

    for page_num in range(1, num_pages + 1):
        # Background
        c.setFillColor(bg_color)
        c.rect(0, 0, width, height, fill=True, stroke=False)

        # Header bar
        c.setFillColor(accent_color)
        c.rect(0, height - 60, width, 60, fill=True, stroke=False)

        # Title
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 28)
        c.drawString(30, height - 42, f"Document {variant} - Page {page_num}")

        # Content area
        c.setFillColor(colors.black)
        c.setFont("Helvetica", 14)

        # Add some content that varies between pages
        y = height - 100
        c.drawString(50, y, f"This is page {page_num} of {num_pages}")

        y -= 30
        c.drawString(50, y, f"Variant: {variant}")

        # Add some boxes that differ slightly between A and B
        y -= 50
        for i in range(5):
            box_x = 50 + i * 150
            box_y = y - 100

            # A has blue boxes, B has orange boxes
            c.setFillColor(accent_color)
            c.rect(box_x, box_y, 120, 80, fill=True, stroke=False)

            c.setFillColor(colors.white)
            c.setFont("Helvetica-Bold", 12)
            c.drawString(box_x + 10, box_y + 50, f"Item {i + 1}")
            c.setFont("Helvetica", 10)
            c.drawString(box_x + 10, box_y + 30, f"Page {page_num}")

        # Add differences for some pages (every 10th page has extra content in B)
        if variant == "B" and page_num % 10 == 0:
            c.setFillColor(colors.HexColor("#FF5722"))
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, 150, f"*** MODIFIED IN VERSION B ***")
            c.rect(50, 50, 200, 80, fill=True, stroke=False)

        # Page number at bottom
        c.setFillColor(colors.gray)
        c.setFont("Helvetica", 10)
        c.drawString(width / 2 - 20, 20, f"Page {page_num} / {num_pages}")

        c.showPage()

        if page_num % 100 == 0:
            print(f"  Created {page_num}/{num_pages} pages...")

    c.save()
    print(f"  Saved: {filename}")


def create_pptx(filename: str, num_slides: int, variant: str = "A"):
    """Create a PowerPoint with specified number of slides.

    Args:
        filename: Output filename
        num_slides: Number of slides to create
        variant: "A" or "B" - affects colors and some content
    """
    print(f"Creating {filename} with {num_slides} slides...")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Different colors for A and B variants
    if variant == "A":
        accent_rgb = RGBColor(25, 118, 210)  # Blue
        bg_rgb = RGBColor(227, 242, 253)  # Light blue
    else:
        accent_rgb = RGBColor(245, 124, 0)  # Orange
        bg_rgb = RGBColor(255, 243, 224)  # Light orange

    for slide_num in range(1, num_slides + 1):
        slide = prs.slides.add_slide(blank_layout)

        # Background shape
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_rgb
        bg.line.fill.background()

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = accent_rgb
        header.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.15),
            Inches(10), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"Document {variant} - Slide {slide_num}"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Content text
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2),
            Inches(12), Inches(1)
        )
        content_frame = content_box.text_frame
        content_para = content_frame.paragraphs[0]
        content_para.text = f"This is slide {slide_num} of {num_slides} | Variant: {variant}"
        content_para.font.size = Pt(14)

        # Add content boxes
        for i in range(5):
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5 + i * 2.5), Inches(2.5),
                Inches(2), Inches(1.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = accent_rgb
            box.line.fill.background()

            # Box label
            box_label = slide.shapes.add_textbox(
                Inches(0.6 + i * 2.5), Inches(2.7),
                Inches(1.8), Inches(1)
            )
            label_frame = box_label.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = f"Item {i + 1}\nSlide {slide_num}"
            label_para.font.size = Pt(12)
            label_para.font.bold = True
            label_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add differences for some slides (every 10th slide has extra content in B)
        if variant == "B" and slide_num % 10 == 0:
            diff_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(5),
                Inches(4), Inches(1.5)
            )
            diff_box.fill.solid()
            diff_box.fill.fore_color.rgb = RGBColor(255, 87, 34)  # Deep orange
            diff_box.line.fill.background()

            diff_label = slide.shapes.add_textbox(
                Inches(0.6), Inches(5.3),
                Inches(3.8), Inches(1)
            )
            diff_frame = diff_label.text_frame
            diff_para = diff_frame.paragraphs[0]
            diff_para.text = "*** MODIFIED IN VERSION B ***"
            diff_para.font.size = Pt(16)
            diff_para.font.bold = True
            diff_para.font.color.rgb = RGBColor(255, 255, 255)

        # Page number at bottom
        page_num_box = slide.shapes.add_textbox(
            Inches(6), Inches(7),
            Inches(2), Inches(0.4)
        )
        page_frame = page_num_box.text_frame
        page_para = page_frame.paragraphs[0]
        page_para.text = f"Slide {slide_num} / {num_slides}"
        page_para.font.size = Pt(10)
        page_para.font.color.rgb = RGBColor(128, 128, 128)
        page_para.alignment = PP_ALIGN.CENTER

        if slide_num % 100 == 0:
            print(f"  Created {slide_num}/{num_slides} slides...")

    prs.save(filename)
    print(f"  Saved: {filename}")


def main():
    """Create all test files."""
    output_dir = Path(__file__).parent
    num_pages = 1000

    print(f"\n=== Creating {num_pages}-page test files ===\n")

    # Create PDFs
    create_pdf(str(output_dir / "A_1000.pdf"), num_pages, "A")
    create_pdf(str(output_dir / "B_1000.pdf"), num_pages, "B")

    # Create PPTXs
    create_pptx(str(output_dir / "A_1000.pptx"), num_pages, "A")
    create_pptx(str(output_dir / "B_1000.pptx"), num_pages, "B")

    print("\n=== All files created! ===")
    print(f"\nFiles created in: {output_dir}")
    print("  - A_1000.pdf (1000 pages)")
    print("  - B_1000.pdf (1000 pages, with differences every 10 pages)")
    print("  - A_1000.pptx (1000 slides)")
    print("  - B_1000.pptx (1000 slides, with differences every 10 slides)")


if __name__ == "__main__":
    main()
